"""
Core engine that turns an Excel workbook of people (one sheet per team)
into a branded PowerPoint deck, one grid-of-cards slide per page.

Expected columns per sheet (case-insensitive, order doesn't matter):
    Name | Phone | Email | Designation | LinkedinUrl | Location | PhotoUrl

Only Name and Designation are required. LinkedinUrl (if present and a real
URL) is used to hyperlink the person's name. PhotoUrl (if present and a real
image) is downloaded and cropped into a circle; anything else falls back to
a generated placeholder avatar.
"""

import io
import re
from dataclasses import dataclass, field

import requests
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------
# Brand / layout constants -- tweak these to re-skin the deck
# ----------------------------------------------------------------------
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

HEADER_COLOR_LEFT = RGBColor(0x6A, 0x0D, 0x83)   # deep purple
HEADER_COLOR_RIGHT = RGBColor(0xC0, 0x00, 0x6E)  # magenta
CARD_FILL = RGBColor(0xFA, 0xD6, 0xD9)           # soft pink
CARD_BORDER = RGBColor(0xE9, 0xA2, 0xA5)
FEATURED_FILL = RGBColor(0xB8, 0x0A, 0x83)       # bold magenta for the "top" card
FEATURED_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
NAME_LINK_COLOR = RGBColor(0x11, 0x55, 0xCC)
DESIG_COLOR = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

GRID_COLS = 3
GRID_ROWS_NORMAL = 3          # 9 per page when there's no featured card
GRID_ROWS_WITH_FEATURED = 2   # 6 per page (+1 featured) when first row is featured
PEOPLE_PER_PAGE_NORMAL = GRID_COLS * GRID_ROWS_NORMAL
PEOPLE_PER_PAGE_FEATURED = GRID_COLS * GRID_ROWS_WITH_FEATURED  # + 1 featured

AVATAR_PX = 300  # source resolution for generated/circular avatars

REQUIRED_COLS = ["name", "designation"]
OPTIONAL_COLS = ["phone", "email", "linkedinurl", "location", "photourl"]


@dataclass
class Person:
    name: str = ""
    designation: str = ""
    linkedin_url: str = ""
    photo_url: str = ""
    location: str = ""
    email: str = ""
    phone: str = ""


@dataclass
class Team:
    sheet_name: str
    people: list = field(default_factory=list)
    featured_first: bool = False  # True => first person gets the big top card


# ----------------------------------------------------------------------
# Excel parsing
# ----------------------------------------------------------------------
def load_teams_from_workbook(file_like, featured_sheets=None) -> list:
    """
    file_like: path or file-like object for the .xlsx
    featured_sheets: iterable of sheet names (case-insensitive) whose first
                      row should render as the big "featured" card
                      (e.g. "Executive Management" -> Secretary General).
    """
    import openpyxl

    featured_sheets = {s.lower().strip() for s in (featured_sheets or [])}
    wb = openpyxl.load_workbook(file_like, data_only=True)

    teams = []
    for ws in wb.worksheets:
        header_row = [
            (c.value or "").strip().lower().replace(" ", "")
            if isinstance(c.value, str) else ""
            for c in next(ws.iter_rows(min_row=1, max_row=1))
        ]
        col_idx = {name: i for i, name in enumerate(header_row)}

        if "name" not in col_idx:
            continue  # not a people sheet, skip quietly

        people = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            def get(col):
                i = col_idx.get(col)
                if i is None or i >= len(row):
                    return ""
                val = row[i]
                return str(val).strip() if val is not None else ""

            name = get("name")
            if not name:
                continue
            people.append(Person(
                name=name,
                designation=get("designation"),
                linkedin_url=get("linkedinurl"),
                photo_url=get("photourl"),
                location=get("location"),
                email=get("email"),
                phone=get("phone"),
            ))

        if people:
            teams.append(Team(
                sheet_name=ws.title,
                people=people,
                featured_first=ws.title.lower().strip() in featured_sheets,
            ))
    return teams


# ----------------------------------------------------------------------
# Image handling
# ----------------------------------------------------------------------
LINKEDIN_HOST_RE = re.compile(r"linkedin\.com", re.I)
_BROWSER_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
        "(KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36"
    ),
    "Accept-Language": "en-US,en;q=0.9",
}


def _looks_like_image_url(url: str) -> bool:
    return bool(url) and url.startswith(("http://", "https://"))


def _download_image(url: str, timeout=8):
    """Try to download `url` directly as an image. Returns a PIL Image or None."""
    if not _looks_like_image_url(url):
        return None
    try:
        resp = requests.get(url, timeout=timeout, headers=_BROWSER_HEADERS)
        resp.raise_for_status()
        img = Image.open(io.BytesIO(resp.content)).convert("RGBA")
        return img
    except Exception:
        return None


def _extract_og_image(page_url: str, timeout=8):
    """
    Fetch a normal web page (e.g. a LinkedIn profile URL) and pull the
    `og:image` meta tag out of its HTML -- that's the same photo LinkedIn
    itself shows in link-preview cards, and it's readable without logging
    in. Returns an image URL string, or None.

    Note: LinkedIn increasingly serves a login-wall page to anonymous
    requests, in which case og:image is LinkedIn's own logo rather than the
    person's photo, or the request may simply be blocked. This is
    best-effort -- when it doesn't work, the caller falls back to the
    placeholder avatar rather than failing.
    """
    if not _looks_like_image_url(page_url):
        return None
    try:
        resp = requests.get(page_url, timeout=timeout, headers=_BROWSER_HEADERS)
        resp.raise_for_status()
        html = resp.text
        m = re.search(
            r'<meta[^>]+property=["\']og:image["\'][^>]+content=["\']([^"\']+)["\']',
            html, re.I,
        )
        if not m:
            m = re.search(
                r'<meta[^>]+content=["\']([^"\']+)["\'][^>]+property=["\']og:image["\']',
                html, re.I,
            )
        if m:
            return m.group(1)
    except Exception:
        pass
    return None


def _resolve_avatar_image(photo_url: str, linkedin_url: str):
    """
    Try, in order:
      1. PhotoUrl as a direct image
      2. If PhotoUrl is itself a linkedin.com link (profile page, not an
         image file), pull its og:image and download that
      3. LinkedinUrl as a direct image (rare, but harmless to try)
      4. LinkedinUrl's og:image
    Returns a PIL Image, or None if nothing worked.
    """
    candidates = [u for u in (photo_url, linkedin_url) if u]
    seen = set()
    for url in candidates:
        if url in seen:
            continue
        seen.add(url)

        img = _download_image(url)
        if img is not None:
            return img

        if LINKEDIN_HOST_RE.search(url):
            og_url = _extract_og_image(url)
            if og_url and og_url not in seen:
                seen.add(og_url)
                img = _download_image(og_url)
                if img is not None:
                    return img
    return None


def _make_placeholder_avatar(size=AVATAR_PX) -> Image.Image:
    """Simple generic silhouette avatar, generated on the fly (no asset file needed)."""
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    draw.ellipse([0, 0, size, size], fill=(210, 210, 215, 255))
    # head
    head_r = size * 0.17
    draw.ellipse([size / 2 - head_r, size * 0.22, size / 2 + head_r, size * 0.22 + 2 * head_r],
                 fill=(150, 150, 158, 255))
    # shoulders
    draw.pieslice([size * 0.12, size * 0.55, size * 0.88, size * 1.35], 180, 360,
                  fill=(150, 150, 158, 255))
    return img


def _circular_crop(img: Image.Image, size=AVATAR_PX) -> Image.Image:
    img = img.convert("RGBA")
    # center-crop to square first
    w, h = img.size
    side = min(w, h)
    left, top = (w - side) // 2, (h - side) // 2
    img = img.crop((left, top, left + side, top + side)).resize((size, size), Image.LANCZOS)

    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse([0, 0, size, size], fill=255)
    out = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    out.paste(img, (0, 0), mask)
    return out


_avatar_cache = {}


def get_avatar_stream(photo_url: str, linkedin_url: str = "") -> io.BytesIO:
    """Returns a PNG BytesIO of a circular avatar -- downloaded photo, or a placeholder.

    Tries `photo_url` as a direct image first; if that fails and either URL
    points at linkedin.com, falls back to extracting the profile's
    og:image (see `_resolve_avatar_image`).
    """
    key = f"{photo_url}|{linkedin_url}"
    if key in _avatar_cache:
        _avatar_cache[key].seek(0)
        return _avatar_cache[key]

    img = _resolve_avatar_image(photo_url, linkedin_url)
    if img is None:
        img = _make_placeholder_avatar()
    circ = _circular_crop(img)

    buf = io.BytesIO()
    circ.save(buf, format="PNG")
    buf.seek(0)
    _avatar_cache[key] = buf
    return buf


# ----------------------------------------------------------------------
# Slide building helpers
# ----------------------------------------------------------------------
def _add_gradient_header(slide, title_text, subtitle_right=None):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85))
    header.line.fill.background()
    header.shadow.inherit = False
    fill = header.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = HEADER_COLOR_LEFT
    stops[0].position = 0.0
    stops[1].color.rgb = HEADER_COLOR_RIGHT
    stops[1].position = 1.0
    fill.gradient_angle = 0.0

    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(7.5), Inches(0.7))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return header


def _add_footer(slide, page_label=""):
    tb = slide.shapes.add_textbox(Inches(0.1), SLIDE_H - Inches(0.3), Inches(3), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Company Restricted"
    r.font.size = Pt(8)
    r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    if page_label:
        tb2 = slide.shapes.add_textbox(SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.3), Inches(1.0), Inches(0.25))
        p2 = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = page_label
        r2.font.size = Pt(8)
        r2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)


def _set_no_shadow(shape):
    shape.shadow.inherit = False


def _add_person_card(slide, person: Person, left, top, width, height, featured=False):
    fill_color = FEATURED_FILL if featured else CARD_FILL
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if featured:
        card.line.fill.background()
    else:
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(0.75)
    _set_no_shadow(card)

    photo_size = Inches(0.55) if not featured else Inches(0.7)
    photo_left = left + Inches(0.12)
    photo_top = top + (height - photo_size) / 2
    avatar_stream = get_avatar_stream(person.photo_url, person.linkedin_url)
    pic = slide.shapes.add_picture(avatar_stream, photo_left, photo_top, width=photo_size, height=photo_size)
    # circular crop via preset geometry so PowerPoint also renders it round
    pic.line.fill.background()

    text_left = photo_left + photo_size + Inches(0.15)
    text_width = width - (text_left - left) - Inches(0.12)

    # Name -- its own separate textbox (not grouped with designation), so it
    # can be selected/moved/edited independently in PowerPoint.
    name_h = Inches(0.32) if featured else Inches(0.26)
    name_top = top + Inches(0.08) if featured else top + Inches(0.06)
    name_box = slide.shapes.add_textbox(text_left, name_top, text_width, name_h)
    name_tf = name_box.text_frame
    name_tf.word_wrap = True
    name_tf.margin_left = 0
    name_tf.margin_right = 0
    name_tf.margin_top = 0
    name_tf.margin_bottom = 0
    p_name = name_tf.paragraphs[0]
    r_name = p_name.add_run()
    r_name.text = person.name or "—"
    r_name.font.bold = True
    r_name.font.size = Pt(12 if featured else 10)
    r_name.font.underline = bool(person.linkedin_url)
    r_name.font.color.rgb = WHITE if featured else NAME_LINK_COLOR
    if person.linkedin_url and person.linkedin_url.startswith("http"):
        r_name.hyperlink.address = person.linkedin_url

    # Designation -- separate textbox below the name, independently editable.
    desig_top = name_top + name_h
    desig_h = height - (desig_top - top) - Inches(0.05)
    desig_box = slide.shapes.add_textbox(text_left, desig_top, text_width, desig_h)
    desig_tf = desig_box.text_frame
    desig_tf.word_wrap = True
    desig_tf.margin_left = 0
    desig_tf.margin_right = 0
    desig_tf.margin_top = 0
    desig_tf.margin_bottom = 0
    p_desig = desig_tf.paragraphs[0]
    r_desig = p_desig.add_run()
    r_desig.text = person.designation or ""
    r_desig.font.size = Pt(10 if featured else 9)
    r_desig.font.color.rgb = WHITE if featured else DESIG_COLOR


def _paginate(people, per_page):
    return [people[i:i + per_page] for i in range(0, len(people), per_page)]


def build_team_slides(prs: Presentation, team: Team, blank_layout):
    margin = Inches(0.12)
    grid_top = Inches(1.05)
    gap = Inches(0.12)

    if team.featured_first and team.people:
        featured_person, rest = team.people[0], team.people[1:]
        per_page = PEOPLE_PER_PAGE_FEATURED
        pages = _paginate(rest, per_page)
        if not pages:
            pages = [[]]
    else:
        featured_person, rest = None, team.people
        per_page = PEOPLE_PER_PAGE_NORMAL
        pages = _paginate(rest, per_page)

    total_pages = len(pages)
    card_w = (SLIDE_W - 2 * margin - (GRID_COLS - 1) * gap) / GRID_COLS
    card_h = Inches(0.85)

    for page_i, page_people in enumerate(pages, start=1):
        slide = prs.slides.add_slide(blank_layout)
        title = team.sheet_name
        if total_pages > 1:
            title = f"{team.sheet_name} ({page_i}/{total_pages})"
        _add_gradient_header(slide, title)

        cur_top = grid_top
        if page_i == 1 and featured_person is not None:
            feat_w = Inches(3.1)
            feat_left = (SLIDE_W - feat_w) / 2
            _add_person_card(slide, featured_person, feat_left, cur_top, feat_w, Inches(0.85), featured=True)
            cur_top += Inches(0.85) + Inches(0.22)

        for idx, person in enumerate(page_people):
            row, col = divmod(idx, GRID_COLS)
            left = margin + col * (card_w + gap)
            top = cur_top + row * (card_h + gap)
            _add_person_card(slide, person, left, top, card_w, card_h)

        _add_footer(slide, page_label=f"{page_i}/{total_pages}" if total_pages > 1 else "")


# ----------------------------------------------------------------------
# Top-level entry point
# ----------------------------------------------------------------------
def generate_presentation(
    teams: list,
    title="Account Intelligence Report",
    subtitle="",
    template_path=None,
) -> io.BytesIO:
    prs = Presentation(template_path) if template_path else Presentation()
    if not template_path:
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    # Title slide
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = HEADER_COLOR_LEFT
    bg.line.fill.background()
    _set_no_shadow(bg)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(7), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = WHITE

    # Agenda-ish / one section per sheet
    for team in teams:
        build_team_slides(prs, team, blank_layout)

    # Closing slide
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = HEADER_COLOR_RIGHT
    bg.line.fill.background()
    _set_no_shadow(bg)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(6), Inches(1))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Thank you"
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = WHITE

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out
